import os
import uuid
import pytesseract
import sys
import time
import io
import fitz  # PyMuPDF
from PIL import Image as PILImage
from PIL import ImageOps, ImageFilter
from pptx import Presentation
from shared_utils_app import get_gemini_client, get_chroma_collection

# Setup Tesseract


def get_tesseract_path():
    # Check if the application is running as a bundled executable
    if getattr(sys, 'frozen', False):
        # Path when running as .exe (PyInstaller's temp folder)
        base_path = sys._MEIPASS
    else:
        # Path when running as a normal .py script
        base_path = os.path.dirname(__file__)
        
    # Join with your specific folder name "Tessereact"
    return os.path.join(base_path, "Tessereact", "tesseract.exe")

# Apply the path to pytesseract
pytesseract.pytesseract.tesseract_cmd = get_tesseract_path()

# Define path for extracted images (for the "Vision" feature)
IMAGE_STORE_DIR = "data/images"
if not os.path.exists(IMAGE_STORE_DIR):
    os.makedirs(IMAGE_STORE_DIR)

def extract_text_with_metadata(file_path):
    """
    Enhanced extraction with Image Pre-processing for better OCR accuracy.
    """
    pages_data = [] 
    base_name = os.path.basename(file_path)
    
    # --- PPTX EXTRACTION ---
    if file_path.endswith('.pptx'):
        try:
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                slide_text = []
                img_count = 0
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                    if shape.shape_type == 6: 
                        for s in shape.shapes:
                            if hasattr(s, "text"): slide_text.append(s.text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                slide_text.append(cell.text)
                    
                    if shape.shape_type == 13: # Picture
                        img_count += 1
                        try:
                            image_bytes = shape.image.blob
                            img = PILImage.open(io.BytesIO(image_bytes))
                            
                            # 1. Save original for UI
                            img_filename = f"{base_name}_slide{i+1}_img{img_count}.png"
                            img.save(os.path.join(IMAGE_STORE_DIR, img_filename))
                            
                            # 2. PRE-PROCESS for OCR
                            ocr_img = img.convert('L')
                            ocr_img = ImageOps.autocontrast(ocr_img) # Boost contrast
                            
                            # Upscale 2x for better character recognition
                            w, h = ocr_img.size
                            ocr_img = ocr_img.resize((w*2, h*2), PILImage.Resampling.LANCZOS)
                            
                            ocr_text = pytesseract.image_to_string(ocr_img)
                            if ocr_text.strip():
                                slide_text.append(f"[Diagram Content]: {ocr_text.strip()}")
                        except: pass
                
                final_text = "\n".join(slide_text).strip()
                if final_text:
                    pages_data.append((final_text, i + 1))
        except Exception as e:
            print(f"PPTX Error: {e}")

    # --- PDF EXTRACTION ---
    elif file_path.endswith('.pdf'):
        try:
            doc = fitz.open(file_path)
            for i, page in enumerate(doc):
                text = page.get_text().strip()
                
                # Image Extraction Logic
                image_list = page.get_images(full=True)
                for img_index, img_info in enumerate(image_list):
                    xref = img_info[0]
                    base_image = doc.extract_image(xref)
                    img = PILImage.open(io.BytesIO(base_image["image"]))
                    img_filename = f"{base_name}_page{i+1}_img{img_index+1}.png"
                    img.save(os.path.join(IMAGE_STORE_DIR, img_filename))

                if not text:
                    pix = page.get_pixmap()
                    # Apply pre-processing to scanned PDF pages
                    img = PILImage.frombytes("RGB", [pix.width, pix.height], pix.samples).convert('L')
                    img = ImageOps.autocontrast(img)
                    text = pytesseract.image_to_string(img).strip()
                
                if text:
                    pages_data.append((text, i + 1))
            doc.close()
        except Exception as e:
            print(f"PDF Error: {e}")

    # --- STANDALONE IMAGE OCR ---
    elif file_path.lower().endswith(('.png', '.jpg', '.jpeg')):
        try:
            img = PILImage.open(file_path)
            img.save(os.path.join(IMAGE_STORE_DIR, base_name))
            
            # Pre-process standalone image
            ocr_img = img.convert('L')
            ocr_img = ImageOps.autocontrast(ocr_img)
            text = pytesseract.image_to_string(ocr_img)
            
            if text.strip():
                pages_data.append((text.strip(), 1))
        except Exception: pass

    elif file_path.endswith('.txt'):
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
            if text.strip():
                pages_data.append((text.strip(), 1))
                
    return pages_data

def get_chunks(text, page_num, source_name, chunk_size=1000, overlap=50):
    chunks = []
    for i in range(0, len(text), chunk_size - overlap):
        chunk_content = text[i:i + chunk_size]
        chunks.append({
            "text": chunk_content,
            "metadata": {
                "source": source_name,
                "page": page_num,
                "id": f"{source_name}_p{page_num}_{i}"
            }
        })
    return chunks

def process_files_to_db(file_paths, api_key, db_path, collection_name, progress_callback=None):
    client = get_gemini_client(api_key)
    collection, _ = get_chroma_collection(client, db_path, collection_name)
    
    all_chunks = []
    files_processed = 0
    
    for i, path in enumerate(file_paths):
        fname = os.path.basename(path)
        if progress_callback:
            progress_callback(f"Analyzing: {fname}...", (i / len(file_paths)) * 0.3)
        
        file_data = extract_text_with_metadata(path)
        if file_data:
            files_processed += 1
            for text, pnum in file_data:
                file_chunks = get_chunks(text, pnum, fname)
                all_chunks.extend(file_chunks)

    if not all_chunks:
        return "No text could be extracted. Check file content."

    batch_size = 15 
    total_chunks = len(all_chunks)
    
    for i in range(0, total_chunks, batch_size):
        batch_data = all_chunks[i:i + batch_size]
        docs = [item["text"] for item in batch_data]
        metas = [item["metadata"] for item in batch_data]
        ids = [str(uuid.uuid4()) for _ in range(len(batch_data))]
        
        progress_percent = 0.3 + ((i / total_chunks) * 0.7)
        if progress_callback:
            progress_callback(f"Embedding: {i}/{total_chunks} chunks...", progress_percent)

        success = False
        while not success:
            try:
                collection.add(documents=docs, ids=ids, metadatas=metas)
                success = True
                time.sleep(2) 
            except Exception as e:
                if "429" in str(e):
                    if progress_callback: progress_callback("Quota Full. Waiting 60s...", progress_percent)
                    time.sleep(65)
                else:
                    return f"Database Error: {e}"

    return f"Success! Added {files_processed} files ({total_chunks} chunks) and saved diagrams. 🚀"

def get_unique_sources(client, db_path, collection_name="university_notes"):
    try:
        collection, _ = get_chroma_collection(client, db_path, collection_name)
        results = collection.get(include=["metadatas"])
        if not results['metadatas']: return []
        sources = {m.get("source") for m in results['metadatas'] if m and "source" in m}
        return sorted(list(sources))
    except Exception: return []

def delete_source_from_db(client, db_path, filename, collection_name="university_notes"):
    try:
        collection, _ = get_chroma_collection(client, db_path, collection_name)
        collection.delete(where={"source": filename})
        return f"Removed {filename} successfully."
    except Exception as e: return f"Delete Error: {e}"